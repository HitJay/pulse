"""
RIC Monthly Report — June 2026 (QYJI)
From-scratch python-pptx builder. Run:
    conda run -n research python /tmp/build_monthly_report_v3.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colours ──
NN_BLUE      = RGBColor(0x00, 0x38, 0x65)
NN_BLUE_LITE = RGBColor(0xCC, 0xD9, 0xE8)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY    = RGBColor(0x26, 0x26, 0x26)
GREY_BG      = RGBColor(0xF2, 0xF2, 0xF2)

# ── Font sizes (user-confirmed 2026-05-28) ──
SZ_HEADER    = 11
SZ_BODY      = 9.5
SZ_LINK      = 8.5
SZ_TITLE_BAR = 14
LINE_SPACING = 12

# ── Helpers ──
def set_cell_bg(cell, rgb):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ['a:solidFill', 'a:gradFill', 'a:noFill']:
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')

def add_run(para, text, bold=False, size_pt=SZ_BODY, color=DARK_GREY, italic=False):
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return run

def fill_cell(cell, lines, header=False, category=False, link=False):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.04)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    first = True
    for (text, bold, italic) in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_before = Pt(0.5)
        para.space_after  = Pt(0.5)
        para.line_spacing = Pt(LINE_SPACING)
        c = WHITE if header else (NN_BLUE if category else DARK_GREY)
        sz = SZ_HEADER if header else (SZ_BODY if not link else SZ_LINK)
        add_run(para, text, bold=bold, italic=italic, size_pt=sz, color=c)

def vmerge_restart(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in tcPr.findall(qn('a:vMerge')):
        tcPr.remove(el)
    vm = etree.SubElement(tcPr, qn('a:vMerge'))
    vm.set('val', 'restart')

def vmerge_cont(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in tcPr.findall(qn('a:vMerge')):
        tcPr.remove(el)
    etree.SubElement(tcPr, qn('a:vMerge'))

J = "https://jira.novonordisk.com/browse/"

# ── Rows ──
ROWS = [
    # 1. Capability — RIC-381
    {
        "cat": "Capability", "cat_span": 3,
        "act": "HepG2 EE — DRUG-seq × imaging cross-modal MoA + vAssay leak-proof benchmark (RIC-381)",
        "desc": [
            ("1. ", False, False),
            ("1440-well three-modality alignment", True, False),
            (" with zero gaps: DRUG-seq transcriptomics × C24 imaging (DINOv2 384-dim + Seahorse prediction) × TMRM mechanism axis. Pipeline + tests in vCell (24/24 passed).", False, False),
            ("2. ", False, False),
            ("vAssay model review — label-leak diagnosis", True, False),
            (": traced the random R²≈0.77 to within-group y-leak (264 imaging rows → only 88 unique Seahorse y; within-group std≈0) + plate effect; established well-group-aware + leave-plate-out evaluation framework.", False, False),
            ("3. ", False, False),
            ("Cross-modal MoA story + EE-DrugSeq-v1 benchmark", True, False),
            (": 175 target perturbations, KD-quality tiering (strong 46 / weak 94 / failed 28); MoA states uncoupler-like 53 / mixed 39 / biogenesis-like 13 / toxic-collapse 6 / neutral 64; ", False, False),
            ("52 tier-1 immediate-validation candidates", True, False),
            (" (DDI2, TAGLN, G6PC, …).", False, False),
            ("4. TMRM channel-identity correction verified against 1_Pipeline/2_tmrm.py and imaging.", False, False),
        ],
        "link": [(f"{J}RIC-381", False, False)],
    },
    # 2. Capability — RIC-349
    {
        "cat": "", "act": "Safety landscape tool — modality-aware + conditional-risk + v6 keywords (RIC-349)",
        "desc": [
            ("1. ", False, False),
            ("Modality-aware prompts + hedge-aware risk scoring", True, False),
            (": peripherally-restricted / partial-IA modality no longer blindly down-weights CNS LoF evidence.", False, False),
            ("2. ", False, False),
            ("P0/P1 scoring + conditional-risk layer", True, False),
            (" (2026-06-18): modality-based mitigation and residual no-go concerns made explicit for safety review.", False, False),
            ("3. Packaged toolkit as installable ", False, False),
            ("safety", False, True),
            (" CLI; produced ZH/EN ", False, False),
            ("one-pager + Mermaid pipeline flowchart", True, False),
            (" (EN PPTX+PDF, 10 slides) for onboarding.", False, False),
            ("4. ", False, False),
            ("v6_composite keyword set", True, False),
            (": binary accuracy ", False, False),
            ("70.8%", True, False),
            (" on NCBN N=305 (+2.3 pp vs v5_final).", False, False),
        ],
        "link": [(f"{J}RIC-349", False, False)],
    },
    # 3. Capability — RIC-388
    {
        "cat": "", "act": "BBB prediction expansion — Phase 2/3/4 + auto-report pipeline (RIC-388)",
        "desc": [
            ("1. ", False, False),
            ("Phase 2 (peptide DL)", True, False),
            (": local ESM-2 peptide-BBB model + ESM-BBB-Pred / DeepB3P; 4 BRP (BRINP2-related) sequences predicted for EJNH (Jason).", False, False),
            ("2. ", False, False),
            ("Phase 3 (benchmarking)", True, False),
            (": cross-tool sensitivity / specificity / MCC; B3clf 12-model committee for SM.", False, False),
            ("3. ", False, False),
            ("Phase 4 — protraction-aware fusion layer", True, False),
            (": lipidated / long-acting peptides no longer scored as backbone false-positives; ", False, False),
            ("end-to-end auto-report pipeline", True, False),
            (" (CLI → HTML + PPTX + LLM narrative).", False, False),
            ("4. ", False, False),
            ("Phase 4+ — modality-adaptive uncertainty", True, False),
            (": SM uses committee-disagreement 90% Wilson interval; peptide retains ESM framing — method-aware, no longer mixed.", False, False),
        ],
        "link": [(f"{J}RIC-388", False, False)],
    },
    # 4. Targets — multi-target
    {
        "cat": "Targets", "cat_span": 1,
        "act": "Multi-target safety assessments — RORA / CDK8-19 / DGKQ family / GRB10 / PAM / EE strict-tox",
        "desc": [
            ("• ", False, False),
            ("RORA", True, False),
            (" (req. MUEW): direction-feasibility for cardiometabolic / DIO — activation vs peripherally-restricted partial inverse agonist, modality-aware run. ", False, False),
            ("RORA tool compounds", True, False),
            (" SR3335 / SR1001 / SR1078 BBB report via B3clf / CNS-MPO + committee uncertainty → ", False, False),
            ("RIC-389 Done", True, False),
            (".", False, False),
            ("• ", False, False),
            ("CDK8 / CDK19 inhibition", True, False),
            (": independent cross-check supporting NCBN/IKEB reproductive-tox termination call on RP0909 (CDK8i, T2D).", False, False),
            ("• ", False, False),
            ("DGKQ family", True, False),
            (" (9 paralogs): inhibition RED=2 / YELLOW=5 / GREEN=2; activation-vs-inhibition two-direction comparison delivered.", False, False),
            ("• Additional batches: ", False, False),
            ("GRB10 inhibition", True, False),
            (", ", False, False),
            ("PAM inhibition", True, False),
            (", ", False, False),
            ("EE strict-tox panel", True, False),
            (", and 2026-06-09 / 06-10 GJSN series — outputs mirrored to /TDE_TV/shared_folder/QYJI/safety/.", False, False),
        ],
        "link": [
            (f"{J}RIC-349", False, False),
            (f"{J}RIC-389", False, False),
            (f"{J}RIC-386", False, False),
        ],
    },
    # 5. Campaign — RIC-261
    {
        "cat": "Campaign", "cat_span": 3,
        "act": "Skeletal Muscle Atrophy assay — new batch + 4-way skeleton benchmark + cross-batch generalisation (RIC-261)",
        "desc": [
            ("1. New batch ", False, False),
            ("IGWU_20260608_1 (PLXN/SEMA siRNA repeat)", True, False),
            (": 240 wells end-to-end through v3.5 decoupled pipeline (no code changes; new wet-lab data).", False, False),
            ("2. ", False, False),
            ("4-way skeleton-analysis benchmark", True, False),
            (": ImageJ / OldPipeline / LegacyApprox / NewPyimagej over 240 images.", False, False),
            ("3. ", False, False),
            ("Cross-batch SSMD validation (3 batches, 720 images)", True, False),
            (": per-image total-length ρ ≥ 0.97, SSMD ρ ≥ 0.986, hit-class concordance ≥ 11/12 (best 20/20).", False, False),
            ("4. Single-file HTML dashboard updated: persistent path + cross-batch volcano (time-period & class colouring, click-to-jump).", False, False),
        ],
        "link": [(f"{J}RIC-261", False, False)],
    },
    # 6. Campaign — RIC-298 / RIC-390
    {
        "cat": "", "act": "Adipocyte Lipolysis Assay — imaging arm split + Round-1 BF analysis (RIC-298 / RIC-390)",
        "desc": [
            ("1. Imaging-side activity log split out into new ticket ", False, False),
            ("RIC-390", True, False),
            ("; RIC-298 kept as parent vAssay capability ticket.", False, False),
            ("2. ", False, False),
            ("Round-1 KDE plate BF analysis", True, False),
            (" (data: Yang Huan / OFGM): CIDEC KD vs NTC, ", False, False),
            ("2,700 images (5 Z × 9 fields × 60 wells)", True, False),
            (", DICOM tiff stack → segmentation → droplet-level quantification end-to-end.", False, False),
        ],
        "link": [
            (f"{J}RIC-298", False, False),
            (f"{J}RIC-390", False, False),
        ],
    },
    # 7. Campaign — RIC-340 closure
    {
        "cat": "", "act": "RadEE — Speakman cohort VBQ pipeline validation closure (RIC-340)",
        "desc": [
            ("Closed RIC-340 with full activity log of the only actively pursued sub-stream: collaboration with ", False, False),
            ("NNRCO Imaging Analytics (Mateusz Florkow et al., Oxford)", True, False),
            (" to validate their VBQ MRI-marker docker pipeline on Dr Speakman's cohort. Status: ", False, False),
            ("Done", True, False),
            (".", False, False),
        ],
        "link": [(f"{J}RIC-340", False, False)],
    },
]

VMERGE_SPANS = [
    (1, 3),  # Capability rows 1-3
    (5, 3),  # Campaign rows 5-7
]

OUT = "reports/ric-monthly/2026-06/monthly_report_2026_06.pptx"

# ── Build ──
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

ML = Inches(0.25); MR = Inches(0.25); MT = Inches(0.15)
W = prs.slide_width; H = prs.slide_height

# Header bar
HDR_H = Inches(0.42)
hdr = slide.shapes.add_shape(1, ML, MT, W - ML - MR, HDR_H)
hdr.fill.solid(); hdr.fill.fore_color.rgb = NN_BLUE; hdr.line.fill.background()
tf = hdr.text_frame; tf.word_wrap = False
tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.04)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
add_run(p, "Summary of activities", bold=True, size_pt=SZ_TITLE_BAR, color=WHITE)
add_run(p, "     \u2502     QYJI     \u2502     Page 1/1     \u2502     2026/06/30     \u2502     ", bold=False, size_pt=9, color=NN_BLUE_LITE)
add_run(p, "Novo Nordisk\u00ae", bold=True, size_pt=10, color=WHITE)

# Table
TBL_T = MT + HDR_H + Inches(0.06)
TBL_H = H - TBL_T - Inches(0.1)
TBL_W = W - ML - MR
N_ROWS = 1 + len(ROWS); N_COLS = 4

tbl_shape = slide.shapes.add_table(N_ROWS, N_COLS, ML, TBL_T, TBL_W, TBL_H)
tbl = tbl_shape.table

col_w = [int(TBL_W * 0.08), int(TBL_W * 0.18), int(TBL_W * 0.59), int(TBL_W * 0.15)]
col_w[-1] = TBL_W - sum(col_w[:-1])
for ci, w in enumerate(col_w):
    tbl.columns[ci].width = w

HDR_ROW_H  = Inches(0.28)
DATA_ROW_H = int((TBL_H - HDR_ROW_H) / max(len(ROWS), 1))
tbl.rows[0].height = HDR_ROW_H
for ri in range(1, N_ROWS):
    tbl.rows[ri].height = DATA_ROW_H

for ci, h in enumerate(["Category", "List of activities", "Brief description", "Relevant links/files"]):
    cell = tbl.cell(0, ci)
    set_cell_bg(cell, NN_BLUE)
    fill_cell(cell, [(h, True, False)], header=True)

for ri, row in enumerate(ROWS):
    dr = ri + 1
    bg = GREY_BG if dr % 2 == 0 else WHITE
    cat_cell = tbl.cell(dr, 0)
    set_cell_bg(cat_cell, NN_BLUE_LITE)
    fill_cell(cat_cell, [(row["cat"], True, False)] if row["cat"] else [("", False, False)], category=bool(row["cat"]))
    act_cell = tbl.cell(dr, 1)
    set_cell_bg(act_cell, bg); fill_cell(act_cell, [(row["act"], True, False)])
    desc_cell = tbl.cell(dr, 2)
    set_cell_bg(desc_cell, bg); fill_cell(desc_cell, row["desc"])
    link_cell = tbl.cell(dr, 3)
    set_cell_bg(link_cell, bg); fill_cell(link_cell, row["link"], link=True)

for (start, span) in VMERGE_SPANS:
    vmerge_restart(tbl.cell(start, 0))
    for i in range(1, span):
        vmerge_cont(tbl.cell(start + i, 0))

prs.save(OUT)
print("Saved:", OUT)
