#!/usr/bin/env python3
"""
build_deck.py – 将 Marp Markdown 转换为 PPTX/PDF 幻灯片。

前置条件:
    npm install -g @marp-team/marp-cli

用法:
    python scripts/build_deck.py weekly 2026-W19          # 周报 → PPTX
    python scripts/build_deck.py weekly 2026-W19 --pdf     # 周报 → PDF
    python scripts/build_deck.py monthly 2026-05           # 月报 → PPTX
"""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

PULSE_ROOT = Path(__file__).resolve().parent.parent
REPORTS_DIR = PULSE_ROOT / "reports"
DECKS_DIR = PULSE_ROOT / "decks"


def check_marp() -> str | None:
    """检查 marp 是否安装，返回路径。"""
    marp_path = shutil.which("marp")
    if marp_path:
        return marp_path
    # 尝试 npx
    npx_path = shutil.which("npx")
    if npx_path:
        return f"{npx_path} marp"
    return None


def build(
    report_type: str,
    identifier: str,
    output_format: str = "pptx",
) -> Path | None:
    """
    将报告 Markdown 转换为幻灯片。

    Args:
        report_type: 'weekly' 或 'monthly'
        identifier: 如 '2026-W19' 或 '2026-05'
        output_format: 'pptx' 或 'pdf'
    """
    # 查找源文件
    src = REPORTS_DIR / report_type / f"{identifier}.md"
    if not src.exists():
        print(f"❌ 源文件不存在: {src}")
        print(f"   请先运行: python scripts/render.py {report_type}")
        return None

    # 输出路径
    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    out_file = DECKS_DIR / f"{identifier}.{output_format}"

    # 检查 marp
    marp_cmd = check_marp()
    if marp_cmd is None:
        print("⚠️  marp-cli 未安装，尝试 fallback 方法...")
        return _fallback_build(src, out_file, output_format)

    # 构建命令
    cmd = f"{marp_cmd} {src} --{output_format} -o {out_file} --allow-local-files"

    print(f"🔨 正在生成 {output_format.upper()}...")
    print(f"   源文件: {src}")
    print(f"   输出到: {out_file}")

    try:
        result = subprocess.run(
            cmd,
            shell=True,
            capture_output=True,
            text=True,
            cwd=str(PULSE_ROOT),
        )
        if result.returncode == 0:
            print(f"✅ {output_format.upper()} 已生成: {out_file}")
            return out_file
        else:
            print(f"❌ marp 执行失败:\n{result.stderr}")
            return _fallback_build(src, out_file, output_format)
    except Exception as e:
        print(f"❌ 执行出错: {e}")
        return _fallback_build(src, out_file, output_format)


def _fallback_build(src: Path, out_file: Path, output_format: str) -> Path | None:
    """当 marp 不可用时，使用 python-pptx 做简易 PPTX 转换。"""
    if output_format != "pptx":
        print("   Fallback 仅支持 PPTX。请安装 marp-cli 以生成 PDF。")
        print("   npm install -g @marp-team/marp-cli")
        return None

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("❌ python-pptx 未安装。请运行: pip install python-pptx")
        return None

    print("   使用 python-pptx fallback 生成简易 PPTX...")

    # 读取 markdown，按 --- 分页
    content = src.read_text(encoding="utf-8")

    # 移除 YAML front matter
    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            content = parts[2]

    # 按 --- 分割为幻灯片
    slides_content = [s.strip() for s in content.split("\n---\n") if s.strip()]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for slide_text in slides_content:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = slide_text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                tf.text = line.lstrip("# ")
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True
            else:
                p = tf.add_paragraph()
                p.text = line.lstrip("- ")
                p.font.size = Pt(16)

    prs.save(str(out_file))
    print(f"✅ PPTX (fallback) 已生成: {out_file}")
    return out_file


# ── CLI ──────────────────────────────────────────────────────────


def main():
    import argparse

    parser = argparse.ArgumentParser(description="Pulse – 幻灯片生成")
    parser.add_argument(
        "type", choices=["weekly", "monthly"], help="报告类型"
    )
    parser.add_argument("id", help="标识符，如 2026-W19 或 2026-05")
    parser.add_argument("--pdf", action="store_true", help="生成 PDF 而非 PPTX")
    args = parser.parse_args()

    fmt = "pdf" if args.pdf else "pptx"
    build(args.type, args.id, fmt)


if __name__ == "__main__":
    main()
